import os
import pdfplumber
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook


def extract_text(filepath: str) -> str:
    ext = os.path.splitext(filepath)[1].lower()

    if ext == ".pdf":
        return extract_pdf(filepath)
    elif ext == ".docx":
        return extract_docx(filepath)
    elif ext == ".pptx":
        return extract_pptx(filepath)
    elif ext == ".xlsx":
        return extract_xlsx(filepath)
    elif ext == ".md":
        return extract_md(filepath)
    else:
        return ""


def extract_pdf(path):
    text = ""
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            text += page.extract_text() or ""
    return text


def extract_docx(path):
    doc = Document(path)
    return "\n".join([para.text for para in doc.paragraphs])


def extract_pptx(path):
    prs = Presentation(path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return "\n".join(text_runs)


def extract_xlsx(path):
    wb = load_workbook(path)
    text = ""
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            line = " ".join([str(cell) for cell in row if cell is not None])
            text += line + "\n"
    return text


def extract_md(path):
    with open(path, "r", encoding="utf-8") as f:
        return f.read()

print("parser.py module loaded successfully.")